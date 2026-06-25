#!/usr/bin/env python3
"""Build a professional hackathon PowerPoint presentation."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import math

# ── Colours ──────────────────────────────────────────────────────────────
BLUE       = RGBColor(0x25, 0x63, 0xEB)
INDIGO     = RGBColor(0x4F, 0x46, 0xE5)
SLATE_900  = RGBColor(0x0F, 0x17, 0x2A)
SLATE_700  = RGBColor(0x33, 0x42, 0x55)
SLATE_500  = RGBColor(0x64, 0x74, 0x8B)
SLATE_400  = RGBColor(0x94, 0xA3, 0xB8)
SLATE_200  = RGBColor(0xE2, 0xE8, 0xF0)
SLATE_100  = RGBColor(0xF1, 0xF5, 0xF9)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
GREEN      = RGBColor(0x05, 0x96, 0x69)
AMBER      = RGBColor(0xD9, 0x77, 0x06)
RED        = RGBColor(0xDC, 0x26, 0x26)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height
CW = Inches(12)                     # content width
CH = Inches(6)                      # content height

# ── Helpers ──────────────────────────────────────────────────────────────

def add_bg(slide, color=WHITE):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, left, top, width, height, fill=None, border=None, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    # Adjust radius if needed
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_circle(slide, left, top, size, fill=None, border=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

def add_arrow_right(slide, left, top, width, height, fill=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_arrow_down(slide, left, top, width, height, fill=BLUE):
    shape = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    return shape

def add_line(slide, x1, y1, x2, y2, color=SLATE_200, width=Pt(1.5)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # MSO_CONNECTOR.STRAIGHT
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_textbox(slide, left, top, width, height, text, font_size=18, color=SLATE_900,
                bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    try:
        tf.paragraphs[0].space_before = Pt(0)
        tf.paragraphs[0].space_after = Pt(0)
    except:
        pass
    return txBox

def add_multiline_textbox(slide, left, top, width, height, lines, font_size=18, color=SLATE_900,
                          bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.alignment = alignment
        p.space_after = Pt(font_size * (line_spacing - 1) * 2)
    return txBox

def add_card(slide, left, top, width, height, title, body="", accent=BLUE, title_size=18, body_size=16):
    """Add a card with left accent bar."""
    card = add_rect(slide, left, top, width, height, fill=WHITE, border=SLATE_200)
    # accent bar
    add_rect(slide, left, top, Inches(0.06), height, fill=accent)
    # title
    add_textbox(slide, left + Inches(0.35), top + Inches(0.3),
                width - Inches(0.7), Inches(0.5), title,
                font_size=title_size, bold=True, color=SLATE_900)
    if body:
        add_textbox(slide, left + Inches(0.35), top + Inches(0.85),
                    width - Inches(0.7), height - Inches(1.2), body,
                    font_size=body_size, color=SLATE_500)
    return card

def add_label(slide, left, top, width, height, text, font_size=11, color=SLATE_500, bold=False,
              bg=SLATE_100, border=SLATE_200):
    """Add a pill label."""
    shape = add_rect(slide, left, top, width, height, fill=bg, border=border)
    add_textbox(slide, left, top, width, height, text, font_size=font_size,
                color=color, bold=bold, alignment=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)
    return shape

# ════════════════════════════════════════════════════════════════════════
# SLIDE 0 — TITLE SLIDE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide)

# Large blue accent block at top-left
add_rect(slide, Inches(0), Inches(0), Inches(7), Inches(0.08), fill=BLUE)

# Title
add_textbox(slide, Inches(1), Inches(1.5), Inches(11), Inches(1.2),
            "Cross-Platform Identity Resolution\n& Agentic AI Ad Engine",
            font_size=40, bold=True, color=SLATE_900)

# Subtitle
add_textbox(slide, Inches(1), Inches(3.2), Inches(11), Inches(0.8),
            "A Unified Customer Data Platform for Privacy-Preserving, AI-Driven Personalization",
            font_size=22, color=SLATE_500)

# Separator line
add_rect(slide, Inches(1), Inches(4.3), Inches(3), Inches(0.04), fill=BLUE)

# Meta info
add_textbox(slide, Inches(1), Inches(4.7), Inches(11), Inches(0.5),
            "Epsilon Hackathon  ·  June 2026",
            font_size=20, color=SLATE_400)

# Right decorative element - large "CDP" watermark
add_textbox(slide, Inches(8.5), Inches(4.5), Inches(4.5), Inches(2.5),
            "CDP",
            font_size=96, bold=True, color=RGBColor(0xF1, 0xF5, 0xF9), alignment=PP_ALIGN.RIGHT)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Good morning, everyone. Thank you for joining us today. "
    "We are excited to present our solution — Cross-Platform Identity Resolution and Agentic AI Ad Engine. "
    "A Customer Data Platform that solves identity fragmentation across platforms "
    "and delivers personalized AI-generated advertisements — all powered by local, privacy-preserving AI. "
    "Let us walk you through the problem, our solution, and what we have built."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TEAM INTRODUCTION
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "Meet the Team",
            font_size=36, bold=True, color=SLATE_900)
add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
            "Built end-to-end with equal contributions across backend, frontend, AI, and infrastructure.",
            font_size=20, color=SLATE_500)

# Team member 1
x1, y1 = Inches(0.7), Inches(1.8)
w_card = Inches(5.5)
h_card = Inches(4.8)

card1 = add_rect(slide, x1, y1, w_card, h_card, fill=WHITE, border=SLATE_200)
# Avatar circle
add_circle(slide, x1 + Inches(0.4), y1 + Inches(0.4), Inches(1.2), fill=BLUE)
add_textbox(slide, x1 + Inches(0.4), y1 + Inches(0.4), Inches(1.2), Inches(1.2),
            "AK", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

add_textbox(slide, x1 + Inches(1.9), y1 + Inches(0.45), Inches(3.2), Inches(0.5),
            "Aman Kumar", font_size=26, bold=True, color=SLATE_900)
add_textbox(slide, x1 + Inches(1.9), y1 + Inches(0.95), Inches(3.2), Inches(0.4),
            "Full-Stack Engineer · AI Integration",
            font_size=18, color=BLUE, bold=True)

# Responsibilities
responsibilities = [
    "⚙️  Backend — FastAPI, Redis, MongoDB integration",
    "🤖  AI Pipeline — Ollama LLM, Qdrant vector search",
    "📊  Analytics — Identity resolution engine, evaluation",
    "☁️  Infrastructure — Docker, Docker Compose, Kafka",
]
for i, r in enumerate(responsibilities):
    add_textbox(slide, x1 + Inches(0.4), y1 + Inches(1.85) + Inches(i * 0.6),
                Inches(4.8), Inches(0.5), r, font_size=18, color=SLATE_700)

# Team member 2
x2 = Inches(6.8)
card2 = add_rect(slide, x2, y1, w_card, h_card, fill=WHITE, border=SLATE_200)
add_circle(slide, x2 + Inches(0.4), y1 + Inches(0.4), Inches(1.2), fill=INDIGO)
add_textbox(slide, x2 + Inches(0.4), y1 + Inches(0.4), Inches(1.2), Inches(1.2),
            "AP", font_size=36, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER,
            anchor=MSO_ANCHOR.MIDDLE)

add_textbox(slide, x2 + Inches(1.9), y1 + Inches(0.45), Inches(3.2), Inches(0.5),
            "Abhishek Pundir", font_size=26, bold=True, color=SLATE_900)
add_textbox(slide, x2 + Inches(1.9), y1 + Inches(0.95), Inches(3.2), Inches(0.4),
            "Full-Stack Engineer · Frontend Lead",
            font_size=18, color=BLUE, bold=True)

responsibilities2 = [
    "🎨  Frontend — React, TypeScript, Vite, Tailwind CSS",
    "📈  Dashboard — Recharts, real-time data visualization",
    "🔗  API Gateway — Axios, React Query, event polling",
    "🏗️  Architecture — Component design, routing, state",
]
for i, r in enumerate(responsibilities2):
    add_textbox(slide, x2 + Inches(0.4), y1 + Inches(1.85) + Inches(i * 0.6),
                Inches(4.8), Inches(0.5), r, font_size=18, color=SLATE_700)

# Tech stack badges at bottom
badge_y = y1 + h_card + Inches(0.3)
techs = ["Python", "FastAPI", "React", "TypeScript", "MongoDB", "Redpanda", "Qdrant", "Ollama", "Docker"]
for i, t in enumerate(techs):
    bx = Inches(0.7) + Inches(i * 1.35)
    add_label(slide, bx, badge_y, Inches(1.2), Inches(0.35), t, font_size=12, bold=True,
              bg=RGBColor(0xEE, 0xF2, 0xFF), border=RGBColor(0xC7, 0xD2, 0xFE))

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Our team consists of two full-stack engineers who built this system end-to-end. "
    "Aman Kumar focused on the backend infrastructure — FastAPI, the identity resolution engine, "
    "AI pipeline integration with Ollama and Qdrant, and containerization. "
    "Abhishek Pundir led the frontend development — React, TypeScript, real-time dashboards, "
    "and the overall user experience. We worked collaboratively across all layers, "
    "and this presentation reflects the combined effort of the entire project."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2 — PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "The Problem",
            font_size=36, bold=True, color=SLATE_900)
add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
            "Fragmented customer identities across devices and platforms",
            font_size=22, color=SLATE_500)

# Device journey — horizontal flow
devices = ["Phone", "Laptop", "Tablet", "Website", "Mobile App", "In-Store"]
dev_icons = ["📱", "💻", "📟", "🌐", "📲", "🏪"]
start_x = Inches(0.7)
dev_y = Inches(1.8)
dev_w = Inches(1.6)
dev_h = Inches(1.4)
gap = Inches(0.15)
arrow_gap = Inches(0.4)

for i, (name, icon) in enumerate(zip(devices, dev_icons)):
    dx = start_x + i * (dev_w + arrow_gap)
    # Device card
    card = add_rect(slide, dx, dev_y, dev_w, dev_h, fill=WHITE, border=SLATE_200)
    add_textbox(slide, dx, dev_y + Inches(0.2), dev_w, Inches(0.7),
                icon, font_size=32, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, dx, dev_y + Inches(0.85), dev_w, Inches(0.4),
                name, font_size=16, color=SLATE_700, bold=True, alignment=PP_ALIGN.CENTER)
    # Arrow between
    if i < len(devices) - 1:
        ax = dx + dev_w + Inches(0.05)
        add_arrow_right(slide, ax, dev_y + Inches(0.55), Inches(0.3), Inches(0.25), fill=SLATE_200)

# Separator line
add_line(slide, Inches(0.7), Inches(3.6), Inches(12.3), Inches(3.6), color=SLATE_200, width=Pt(1))

# "Each creates a separate identity" text
add_textbox(slide, Inches(0.7), Inches(3.8), Inches(5), Inches(0.5),
            "Each touchpoint creates a separate, unlinked identity",
            font_size=20, bold=True, color=RED)

# Problem impacts — 4 cards
impacts = [
    ("👤", "Duplicate\nProfiles", "Same customer appears\nas multiple identities"),
    ("🎯", "Poor\nPersonalization", "No unified view for\nrelevant recommendations"),
    ("📉", "Low Marketing\nROI", "Wasted ad spend on\nfragmented audiences"),
    ("🍪", "Cookie\nDeprecation", "Third-party identifiers\nare disappearing"),
]

for i, (emoji, title, desc) in enumerate(impacts):
    cx = Inches(0.7) + Inches(i * 3.05)
    cy = Inches(4.4)
    cw = Inches(2.7)
    ch = Inches(2.5)
    add_rect(slide, cx, cy, cw, ch, fill=WHITE, border=SLATE_200)
    add_textbox(slide, cx, cy + Inches(0.25), cw, Inches(0.5),
                emoji, font_size=28, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx, cy + Inches(0.75), cw, Inches(0.7),
                title, font_size=18, bold=True, color=SLATE_900, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, cx, cy + Inches(1.55), cw, Inches(0.7),
                desc, font_size=14, color=SLATE_500, alignment=PP_ALIGN.CENTER)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Here is the core problem we set out to solve. "
    "A customer browses on their phone, switches to a laptop, visits the website on a tablet, "
    "uses your mobile app, and walks into a store. Every single touchpoint creates a separate, isolated identity. "
    "Businesses end up with duplicate profiles, poor personalization, wasted marketing spend on fragmented audiences, "
    "and with third-party cookies going away, the challenge is only getting worse. "
    "This identity fragmentation is the root cause of poor customer experiences and inefficient marketing."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3 — BUSINESS CONTEXT
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "Business Context",
            font_size=36, bold=True, color=SLATE_900)
add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
            "Why enterprises need a modern Customer Data Platform",
            font_size=22, color=SLATE_500)

# Left: Pain Points → Market → Impact flow
# Pain Point box
pp_x = Inches(0.7)
pp_y = Inches(1.8)
pp_w = Inches(3.2)

add_rect(slide, pp_x, pp_y, pp_w, Inches(4.5), fill=WHITE, border=SLATE_200)
add_rect(slide, pp_x, pp_y, pp_w, Inches(0.06), fill=RED)

add_textbox(slide, pp_x + Inches(0.3), pp_y + Inches(0.25), pp_w - Inches(0.6), Inches(0.5),
            "Business Pain Points", font_size=20, bold=True, color=RED)

pain_points = [
    "•  Customer identities scattered across silos",
    "•  No single view of the customer journey",
    "•  Marketing teams blind to cross-platform behavior",
    "•  Manual identity stitching — error-prone & slow",
    "•  AI models trained on incomplete data",
]
for i, pt in enumerate(pain_points):
    add_textbox(slide, pp_x + Inches(0.3), pp_y + Inches(0.9) + Inches(i * 0.6),
                pp_w - Inches(0.6), Inches(0.5), pt, font_size=16, color=SLATE_700)

# Arrow right
add_arrow_right(slide, pp_x + pp_w + Inches(0.1), Inches(3.6), Inches(0.4), Inches(0.35), fill=BLUE)

# Market opportunity
mo_x = pp_x + pp_w + Inches(0.7)
mo_y = pp_y
mo_w = Inches(3.5)

add_rect(slide, mo_x, mo_y, mo_w, Inches(4.5), fill=WHITE, border=SLATE_200)
add_rect(slide, mo_x, mo_y, mo_w, Inches(0.06), fill=BLUE)

add_textbox(slide, mo_x + Inches(0.3), mo_y + Inches(0.25), mo_w - Inches(0.6), Inches(0.5),
            "Market Opportunity", font_size=20, bold=True, color=BLUE)

market_items = [
    "•  CDP market projected to reach $28B by 2028",
    "•  73% of consumers expect personalized experiences",
    "•  Retail & E-commerce: biggest adoption drivers",
    "•  Cookie deprecation creates urgent need",
    "•  AI-powered CDPs: fastest growing segment",
]
for i, mi in enumerate(market_items):
    add_textbox(slide, mo_x + Inches(0.3), mo_y + Inches(0.9) + Inches(i * 0.6),
                mo_w - Inches(0.6), Inches(0.5), mi, font_size=16, color=SLATE_700)

# Arrow right
add_arrow_right(slide, mo_x + mo_w + Inches(0.1), Inches(3.6), Inches(0.4), Inches(0.35), fill=GREEN)

# Business Impact
bi_x = mo_x + mo_w + Inches(0.7)
bi_y = pp_y
bi_w = Inches(2.8)

add_rect(slide, bi_x, bi_y, bi_w, Inches(4.5), fill=WHITE, border=SLATE_200)
add_rect(slide, bi_x, bi_y, bi_w, Inches(0.06), fill=GREEN)

add_textbox(slide, bi_x + Inches(0.3), bi_y + Inches(0.25), bi_w - Inches(0.6), Inches(0.5),
            "Business Impact", font_size=20, bold=True, color=GREEN)

impact_items = [
    "•  Up to 30% increase in marketing ROI",
    "•  2x improvement in recommendation accuracy",
    "•  Real-time unified customer profiles",
    "•  Privacy-compliant architecture",
    "•  Scalable to millions of customers",
]
for i, ii in enumerate(impact_items):
    add_textbox(slide, bi_x + Inches(0.3), bi_y + Inches(0.9) + Inches(i * 0.6),
                bi_w - Inches(0.6), Inches(0.5), ii, font_size=16, color=SLATE_700)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Enterprises today face a critical challenge: customer identities are scattered across silos — "
    "websites, mobile apps, in-store systems, and marketing platforms. "
    "This fragmentation means no single view of the customer journey, marketing teams operating blind, "
    "and AI models trained on incomplete data. "
    "The Customer Data Platform market is projected to reach 28 billion dollars by 2028, "
    "driven by the urgent need for unified customer profiles, especially with third-party cookie deprecation. "
    "Our solution addresses this market need with an AI-powered CDP that delivers up to 30% improvement in marketing ROI "
    "and real-time unified customer profiles, all while maintaining privacy compliance."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4 — SOLUTION OVERVIEW
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "Solution Overview",
            font_size=36, bold=True, color=SLATE_900)
add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
            "End-to-end pipeline from data ingestion to AI-generated advertisements",
            font_size=22, color=SLATE_500)

# Pipeline flow
steps = [
    ("📊", "Synthetic\nData", "Multi-platform\nclickstream generator"),
    ("📨", "Event\nStreaming", "Redpanda Kafka\nmessage broker"),
    ("🔗", "Identity\nResolution", "Deterministic +\nprobabilistic matching"),
    ("👤", "Unified\nProfile", "Single customer\nview in MongoDB"),
    ("🧠", "Intent\nProfiling", "Ollama LLM\nintent analysis"),
    ("🎯", "Product\nMatching", "Qdrant vector\nsemantic search"),
    ("📣", "AI\nAdvertisement", "Personalized ad\ncreative generation"),
    ("📊", "Real-time\nDashboard", "React + Recharts\nvisualization"),
]

step_w = Inches(1.3)
step_h = Inches(2.6)
start_x = Inches(0.5)
step_y = Inches(1.7)
step_gap = Inches(0.1)
arrow_h = Inches(0.25)

for i, (icon, title, desc) in enumerate(steps):
    sx = start_x + i * (step_w + step_gap)
    # Card
    add_rect(slide, sx, step_y, step_w, step_h, fill=WHITE, border=SLATE_200)
    # Icon
    add_textbox(slide, sx, step_y + Inches(0.2), step_w, Inches(0.5),
                icon, font_size=28, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, sx + Inches(0.1), step_y + Inches(0.7), step_w - Inches(0.2), Inches(0.7),
                title, font_size=14, bold=True, color=SLATE_900, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, sx + Inches(0.1), step_y + Inches(1.55), step_w - Inches(0.2), Inches(0.7),
                desc, font_size=11, color=SLATE_500, alignment=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < len(steps) - 1:
        ax = sx + step_w + Inches(0.02)
        add_arrow_right(slide, ax, step_y + step_h - Inches(1.8), Inches(0.08), arrow_h, fill=BLUE)

# Bottom accent bar
add_rect(slide, Inches(0.5), Inches(4.7), Inches(12.3), Inches(0.04), fill=BLUE)

# Key differentiators at bottom
add_textbox(slide, Inches(0.7), Inches(5.0), Inches(11), Inches(0.5),
            "Key Differentiators",
            font_size=22, bold=True, color=SLATE_900)

diffs = [
    ("🔒", "Privacy-First", "Local LLM, no data leaves\nyour infrastructure"),
    ("⚡", "Real-Time", "Event-driven pipeline\nwith sub-second latency"),
    ("📈", "Scalable", "Kafka-based architecture\nhandles millions of events"),
    ("🤖", "AI-Native", "LLM + Vector Search\nfor deep personalization"),
]

for i, (emoji, title, desc) in enumerate(diffs):
    dx = Inches(0.7) + Inches(i * 3.05)
    dy = Inches(5.5)
    dw = Inches(2.7)
    add_textbox(slide, dx, dy, Inches(0.4), Inches(0.4), emoji, font_size=22)
    add_textbox(slide, dx + Inches(0.5), dy, dw - Inches(0.5), Inches(0.4),
                title, font_size=18, bold=True, color=SLATE_900)
    add_textbox(slide, dx + Inches(0.5), dy + Inches(0.4), dw - Inches(0.5), Inches(0.6),
                desc, font_size=13, color=SLATE_500)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Here is our complete solution at a glance. "
    "Data flows from multi-platform clickstream generators through Redpanda event streaming into our identity resolution engine, "
    "which uses both deterministic and probabilistic matching to build unified customer profiles stored in MongoDB. "
    "From there, our AI pipeline takes over: Ollama LLM analyzes purchase intent, Qdrant vector search finds relevant products, "
    "and the ad generator creates personalized advertisements — all displayed through a real-time React dashboard. "
    "Four key differentiators set us apart: privacy-first architecture with a local LLM so no data ever leaves your infrastructure, "
    "real-time event-driven processing, a scalable Kafka-based pipeline, and deep AI-native personalization."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5 — TECHNICAL ARCHITECTURE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.3), Inches(11), Inches(0.6),
            "Technical Architecture",
            font_size=36, bold=True, color=SLATE_900)

# ── Architecture diagram ──
# We'll create layers as horizontal bands

layers = [
    ("Presentation Layer",  SLATE_100, BLUE, [
        ("React SPA", "Vite + TypeScript\nRecharts + Tailwind"),
        ("Real-time Dashboard", "Live Feed, Analytics\nIdentity Explorer"),
        ("API Client", "Axios + React Query\nEvent polling"),
    ]),
    ("API Layer", RGBColor(0xEE, 0xF2, 0xFF), RGBColor(0x4F, 0x46, 0xE5), [
        ("FastAPI Server", "REST endpoints\nHealth, Events, Profiles"),
        ("Data Validation", "Pydantic models\nRequest/Response schemas"),
        ("API Proxy", "Vite dev proxy\nDocker networking"),
    ]),
    ("AI Layer", RGBColor(0xEC, 0xFD, 0xF5), GREEN, [
        ("Ollama LLM", "Local LLM inference\nIntent profiling"),
        ("Qdrant Vector DB", "Semantic search\nProduct matching"),
        ("Ad Generator", "Agentic AI pipeline\nCreative generation"),
    ]),
    ("Storage Layer", RGBColor(0xFF, 0xFB, 0xEB), AMBER, [
        ("MongoDB", "Unified profiles\nSession storage"),
        ("Redis Cache", "Ad response caching\nPerformance optimization"),
    ]),
    ("Streaming Layer", RGBColor(0xF5, 0xF3, 0xFF), INDIGO, [
        ("Redpanda", "Kafka-compatible\nEvent streaming"),
        ("Event Consumer", "Async processing\nReal-time ingestion"),
    ]),
    ("Data Layer", RGBColor(0xFE, 0xF2, 0xF2), RED, [
        ("Clickstream Generator", "Multi-platform\nsynthetic events"),
        ("Platform Producers", "Platform A & B\nSession data"),
    ]),
]

layer_h = Inches(0.85)
total_layer_h = len(layers) * layer_h
start_y = Inches(1.1)
layer_x = Inches(0.5)
layer_w = Inches(8.5)
legend_x = Inches(9.3)
legend_w = Inches(3.5)

for i, (name, bg_color, accent, items) in enumerate(layers):
    ly = start_y + i * layer_h
    # Layer band background
    add_rect(slide, layer_x, ly, layer_w, layer_h - Inches(0.05), fill=bg_color, border=accent)
    # Layer label on the left
    add_textbox(slide, layer_x + Inches(0.2), ly, Inches(1.8), layer_h - Inches(0.05),
                name, font_size=13, bold=True, color=accent, anchor=MSO_ANCHOR.MIDDLE)
    # Separator
    if i < len(layers) - 1:
        pass  # small gap

    # Items within the layer
    for j, (item_name, item_desc) in enumerate(items):
        ix = layer_x + Inches(2.2) + Inches(j * 2.1)
        iw = Inches(1.95)
        add_textbox(slide, ix, ly + Inches(0.08), iw, Inches(0.35),
                    item_name, font_size=11, bold=True, color=SLATE_900)
        add_textbox(slide, ix, ly + Inches(0.38), iw, Inches(0.4),
                    item_desc, font_size=9, color=SLATE_500)

# ── Legend ──
add_rect(slide, legend_x, start_y, legend_w, total_layer_h, fill=WHITE, border=SLATE_200)
add_textbox(slide, legend_x + Inches(0.2), start_y + Inches(0.15), legend_w - Inches(0.4), Inches(0.4),
            "Architecture Layers", font_size=16, bold=True, color=SLATE_900)

for i, (name, bg_color, accent, _) in enumerate(layers):
    ly = start_y + Inches(0.6) + Inches(i * 0.55)
    # Color dot
    add_rect(slide, legend_x + Inches(0.25), ly, Inches(0.25), Inches(0.25), fill=accent)
    add_textbox(slide, legend_x + Inches(0.65), ly, legend_w - Inches(1), Inches(0.3),
                name, font_size=12, color=SLATE_700, anchor=MSO_ANCHOR.MIDDLE)

# Data flow arrows between layers on the right side of layers
flow_x = layer_x + layer_w + Inches(0.1)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "This is the heart of our system — the technical architecture. "
    "Starting from the bottom, our Data Layer generates multi-platform synthetic clickstream events "
    "that flow into the Streaming Layer powered by Redpanda, a Kafka-compatible message broker. "
    "The Storage Layer uses MongoDB for unified profile persistence and Redis for caching ad responses. "
    "Our AI Layer is where the intelligence lives — Ollama LLM for intent profiling, "
    "Qdrant vector database for semantic product matching, and an agentic pipeline for ad creative generation. "
    "The API Layer built on FastAPI connects everything, and the Presentation Layer renders the React dashboard "
    "with real-time visualizations. "
    "Every layer is containerized with Docker, and the entire system runs with a single Docker Compose command."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6 — PROTOTYPE DEMONSTRATION
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "Prototype Demonstration",
            font_size=36, bold=True, color=SLATE_900)
add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
            "From live event ingestion to AI-generated advertisements",
            font_size=22, color=SLATE_500)

# Screenshot placeholders in a horizontal flow
screens = [
    ("📡", "Live Feed", "Real-time clickstream\nevent ingestion table\nwith filtering & simulation"),
    ("🔍", "Identity Explorer", "Search by Global UID\nview session graph\nmatch method stats"),
    ("👤", "Unified Profile", "Device fingerprints\nlocation history\nintent analysis"),
    ("🧠", "Intent Intelligence", "LLM-generated purchase\nintent profile with\nconfidence scoring"),
    ("🎯", "Product Match", "Qdrant vector search\nsemantic product\nrecommendations"),
    ("📣", "Ad Creative", "Personalized headline\nCTA + product links\nwith refresh capability"),
]

ss_x = Inches(0.35)
ss_y = Inches(1.7)
ss_w = Inches(1.95)
ss_h = Inches(2.8)
ss_gap = Inches(0.08)

for i, (icon, title, desc) in enumerate(screens):
    sx = ss_x + i * (ss_w + ss_gap)
    # Screen card
    add_rect(slide, sx, ss_y, ss_w, ss_h, fill=WHITE, border=SLATE_200)
    # Top accent
    add_rect(slide, sx, ss_y, ss_w, Inches(0.06), fill=BLUE)
    # Icon placeholder
    add_rect(slide, sx + Inches(0.35), ss_y + Inches(0.25), Inches(1.25), Inches(0.9),
             fill=SLATE_100, border=SLATE_200)
    add_textbox(slide, sx + Inches(0.35), ss_y + Inches(0.25), Inches(1.25), Inches(0.9),
                icon, font_size=36, alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # Title
    add_textbox(slide, sx + Inches(0.15), ss_y + Inches(1.3), ss_w - Inches(0.3), Inches(0.4),
                title, font_size=14, bold=True, color=SLATE_900, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, sx + Inches(0.15), ss_y + Inches(1.75), ss_w - Inches(0.3), Inches(0.9),
                desc, font_size=11, color=SLATE_500, alignment=PP_ALIGN.CENTER)

    # Arrow between
    if i < len(screens) - 1:
        ax = sx + ss_w + Inches(0.01)
        add_arrow_right(slide, ax, ss_y + ss_h - Inches(1.8), Inches(0.06), Inches(0.2), fill=BLUE)

# Bottom: app highlights
add_rect(slide, Inches(0.35), Inches(4.8), Inches(12.3), Inches(0.04), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(5.1), Inches(11), Inches(0.5),
            "Application Highlights",
            font_size=22, bold=True, color=SLATE_900)

highlights = [
    "📊  Real-time event streaming with 2-second polling",
    "🔗  Identity resolution with deterministic & probabilistic matching",
    "🧠  Local LLM (Ollama) for privacy-preserving intent profiling",
    "🎯  Vector search (Qdrant) for semantic product recommendations",
    "📣  AI-generated ad creatives with one-click regeneration",
    "📈  Analytics dashboard with F1 scores, precision, recall metrics",
]
for i, h in enumerate(highlights):
    col = i % 2
    row = i // 2
    hx = Inches(0.7) + Inches(col * 6)
    hy = Inches(5.6) + Inches(row * 0.45)
    add_textbox(slide, hx, hy, Inches(5.5), Inches(0.4),
                h, font_size=15, color=SLATE_700)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Let me walk you through the actual application we built. "
    "Starting with the Live Feed — you can see real-time clickstream events being ingested from multiple platforms "
    "with filtering capabilities and an event simulator for testing. "
    "The Identity Explorer lets you search by Global UID and visualize session graphs. "
    "Each profile shows devices, locations, and an LLM-generated intent analysis. "
    "The Ad Studio generates personalized advertisements using our AI pipeline — "
    "the LLM analyzes purchase intent, Qdrant finds relevant products, and the system creates tailored ad copy "
    "with headlines, body text, and CTAs. You can regenerate ads with one click. "
    "The Analytics dashboard provides comprehensive metrics including F1 scores, precision, recall, "
    "platform splits, and infrastructure health monitoring."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AI CAPABILITIES
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "AI Capabilities",
            font_size=36, bold=True, color=SLATE_900)
add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
            "Five AI systems working together to deliver intelligent personalization",
            font_size=22, color=SLATE_500)

# AI cards in a grid
ai_components = [
    ("🎯", "Deterministic\nMatching", "Rule-based identity\nlinking using hashed\nemail, device ID, session",
     BLUE, RGBColor(0xEE, 0xF2, 0xFF)),
    ("🔀", "Probabilistic\nMatching", "ML-based similarity\nscoring across device\nfingerprints & behavior",
     INDIGO, RGBColor(0xF5, 0xF3, 0xFF)),
    ("🧠", "Intent Profiling\n(Ollama)", "Local LLM analyzes\npurchase patterns &\nbehavior for intent",
     GREEN, RGBColor(0xEC, 0xFD, 0xF5)),
    ("🎯", "Vector Search\n(Qdrant)", "Semantic product\nsearch using vector\nembeddings & similarity",
     AMBER, RGBColor(0xFF, 0xFB, 0xEB)),
    ("📣", "Personalized Ad\nGeneration", "Agentic pipeline creates\ntailored ad copy\nheadlines & CTAs",
     RGBColor(0x7C, 0x3A, 0xED), RGBColor(0xF5, 0xF3, 0xFF)),
]

ai_x_start = Inches(0.5)
ai_y = Inches(1.6)
ai_w = Inches(2.35)
ai_h = Inches(3.0)
ai_gap = Inches(0.15)

for i, (icon, title, desc, accent, bg) in enumerate(ai_components):
    ax = ai_x_start + i * (ai_w + ai_gap)
    # Card
    add_rect(slide, ax, ai_y, ai_w, ai_h, fill=bg, border=accent)
    # Top accent
    add_rect(slide, ax, ai_y, ai_w, Inches(0.06), fill=accent)
    # Icon
    add_textbox(slide, ax, ai_y + Inches(0.25), ai_w, Inches(0.5),
                icon, font_size=32, alignment=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, ax + Inches(0.15), ai_y + Inches(0.8), ai_w - Inches(0.3), Inches(0.7),
                title, font_size=16, bold=True, color=SLATE_900, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, ax + Inches(0.15), ai_y + Inches(1.6), ai_w - Inches(0.3), Inches(1.0),
                desc, font_size=13, color=SLATE_500, alignment=PP_ALIGN.CENTER)

    # Arrow between cards
    if i < len(ai_components) - 1:
        arrow_x = ax + ai_w + Inches(0.02)
        add_arrow_right(slide, arrow_x, ai_y + ai_h - Inches(1.6), Inches(0.1), Inches(0.22), fill=SLATE_200)

# Bottom: Key architecture points
add_rect(slide, Inches(0.5), Inches(5.0), Inches(12.3), Inches(0.04), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(5.3), Inches(11), Inches(0.5),
            "Architecture Highlights",
            font_size=22, bold=True, color=SLATE_900)

arch_points = [
    ("🔒", "Local AI", "All models run locally via Ollama\n— zero third-party API calls"),
    ("⚡", "Event-Driven", "Async processing pipeline\nwith Kafka message broker"),
    ("📈", "Scalable", "Horizontally scalable\nmicroservices architecture"),
    ("🛡️", "Privacy First", "Customer data never leaves\nyour infrastructure"),
]

for i, (icon, title, desc) in enumerate(arch_points):
    px = Inches(0.7) + Inches(i * 3.05)
    py = Inches(5.9)
    add_textbox(slide, px, py, Inches(0.4), Inches(0.4), icon, font_size=22)
    add_textbox(slide, px + Inches(0.5), py, Inches(2.3), Inches(0.4),
                title, font_size=18, bold=True, color=SLATE_900)
    add_textbox(slide, px + Inches(0.5), py + Inches(0.4), Inches(2.3), Inches(0.5),
                desc, font_size=13, color=SLATE_500)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Our AI capabilities are built around five core components. "
    "First, deterministic matching uses rule-based identity linking through hashed emails, device IDs, and sessions. "
    "Second, probabilistic matching applies ML-based similarity scoring across device fingerprints and behavioral patterns. "
    "Third, Ollama LLM runs locally to analyze purchase patterns and generate intent profiles — "
    "completely privacy-preserving with no data leaving your infrastructure. "
    "Fourth, Qdrant vector database enables semantic product search using embeddings. "
    "Fifth, our agentic ad generation pipeline creates personalized ad copy with headlines, body text, and CTAs. "
    "The key architectural advantages are all AI runs locally, the entire pipeline is event-driven through Kafka, "
    "microservices scale horizontally, and customer data never leaves your infrastructure."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8 — BUSINESS VALUE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "Business Value",
            font_size=36, bold=True, color=SLATE_900)
add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
            "Tangible outcomes for enterprise customers",
            font_size=22, color=SLATE_500)

# KPI Cards in a grid
kpis = [
    ("👤", "Unified Customer\nProfiles", "Single view of every\ncustomer across all\nplatforms and devices",
     BLUE),
    ("📈", "Higher Marketing\nROI", "Up to 30% improvement\nin ad spend efficiency\nwith targeted campaigns",
     GREEN),
    ("🎯", "Better Customer\nInsights", "Deep intent profiling\nand behavioral\nunderstanding at scale",
     BLUE),
    ("🔒", "Privacy-Centric\nArchitecture", "Local AI inference\nGDPR-ready design\nno data exfiltration",
     INDIGO),
    ("💰", "Reduced Customer\nAcquisition Cost", "Personalized marketing\nreduces CAC by\neliminating waste",
     GREEN),
    ("📊", "Personalized\nMarketing", "AI-generated creatives\ntailored to individual\ncustomer preferences",
     BLUE),
    ("⚡", "Scalable Event\nPipeline", "Kafka-based streaming\nhandles millions of\nevents per second",
     AMBER),
    ("🏗️", "Enterprise Ready\nArchitecture", "Containerized, cloud-\nagnostic deployment\nwith Docker Compose",
     INDIGO),
]

kpi_x_start = Inches(0.5)
kpi_y = Inches(1.6)
kpi_w = Inches(2.9)
kpi_h = Inches(2.3)
kpi_gap_x = Inches(0.15)
kpi_gap_y = Inches(0.15)

for i, (icon, title, desc, accent) in enumerate(kpis):
    col = i % 4
    row = i // 4
    kx = kpi_x_start + col * (kpi_w + kpi_gap_x)
    ky = kpi_y + row * (kpi_h + kpi_gap_y)
    # Card
    add_rect(slide, kx, ky, kpi_w, kpi_h, fill=WHITE, border=SLATE_200)
    # Left accent
    add_rect(slide, kx, ky, Inches(0.05), kpi_h, fill=accent)
    # Icon
    add_textbox(slide, kx + Inches(0.25), ky + Inches(0.25), Inches(0.4), Inches(0.4),
                icon, font_size=24)
    # Title
    add_textbox(slide, kx + Inches(0.7), ky + Inches(0.2), kpi_w - Inches(1), Inches(0.6),
                title, font_size=16, bold=True, color=SLATE_900)
    # Description
    add_textbox(slide, kx + Inches(0.25), ky + Inches(0.9), kpi_w - Inches(0.5), Inches(1.2),
                desc, font_size=13, color=SLATE_500)

# Bottom: key metrics highlight
add_rect(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.04), fill=BLUE)

metrics_text = "30% ↑ Marketing ROI  ·  2x Recommendation Accuracy  ·  Real-time Profiles  ·  Zero Data Leakage  ·  Sub-second Latency"
add_textbox(slide, Inches(0.7), Inches(6.4), Inches(11.5), Inches(0.5),
            metrics_text, font_size=18, bold=True, color=BLUE, alignment=PP_ALIGN.CENTER)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Let us talk about business value. Our solution delivers eight key outcomes for enterprise customers. "
    "Unified customer profiles provide a single view across every platform and device. "
    "Higher marketing ROI with up to 30% improvement in ad spend efficiency through targeted, AI-driven campaigns. "
    "Better customer insights through deep intent profiling and behavioral analysis at scale. "
    "A privacy-centric architecture with local AI inference that is GDPR-ready — no data ever leaves your infrastructure. "
    "Reduced customer acquisition cost through personalized marketing that eliminates waste. "
    "AI-generated ad creatives tailored to individual preferences. "
    "A scalable event pipeline built on Kafka that handles millions of events per second. "
    "And an enterprise-ready, containerized architecture that deploys anywhere with a single command."
)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9 — FUTURE SCOPE
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide)
add_rect(slide, Inches(0), Inches(0), W, Inches(0.06), fill=BLUE)

add_textbox(slide, Inches(0.7), Inches(0.4), Inches(11), Inches(0.7),
            "Future Roadmap",
            font_size=36, bold=True, color=SLATE_900)
add_textbox(slide, Inches(0.7), Inches(1.0), Inches(11), Inches(0.4),
            "Eight-phase roadmap from enterprise hardening to advanced AI capabilities",
            font_size=22, color=SLATE_500)

# Timeline
phases = [
    ("Phase 1", "Authentication\n& Auth", "OAuth 2.0, SSO\nJWT-based security", "Q3 2026", BLUE),
    ("Phase 2", "Real-time\nAnalytics", "Apache Flink streaming\nadvanced aggregations", "Q3 2026", RGBColor(0x7C, 0x3A, 0xED)),
    ("Phase 3", "GDPR APIs", "Data deletion, export\nconsent management", "Q4 2026", GREEN),
    ("Phase 4", "Multi-Tenant\nSaaS", "Organization isolation\nusage-based billing", "Q4 2026", AMBER),
    ("Phase 5", "Cloud\nDeployment", "AWS/GCP/Azure\nKubernetes, Terraform", "Q1 2027", BLUE),
    ("Phase 6", "Predictive\nCLV", "Customer lifetime value\nML forecasting models", "Q1 2027", INDIGO),
    ("Phase 7", "Marketing\nAutomation", "Journey builder\ncampaign triggers", "Q2 2027", GREEN),
    ("Phase 8", "A/B Testing\nEngine", "Multi-variant testing\nstatistical significance", "Q2 2027", AMBER),
]

# Timeline line
timeline_y = Inches(2.3)
add_rect(slide, Inches(0.5), timeline_y, Inches(12.3), Inches(0.04), fill=BLUE)

phase_w = Inches(1.35)
phase_gap = Inches(0.15)
phase_start_x = Inches(0.4)
phase_circle_size = Inches(0.28)

for i, (phase, title, desc, date, accent) in enumerate(phases):
    px = phase_start_x + i * (phase_w + phase_gap)

    # Circle on timeline
    circle_y = timeline_y - phase_circle_size / 2 + Inches(0.02)
    add_circle(slide, px + phase_w / 2 - phase_circle_size / 2, circle_y,
               phase_circle_size, fill=accent)

    # Phase number above
    add_textbox(slide, px, timeline_y - Inches(0.55), phase_w, Inches(0.3),
                phase, font_size=11, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

    # Card below timeline
    card_y = timeline_y + Inches(0.4)
    card_h = Inches(2.2)
    add_rect(slide, px, card_y, phase_w, card_h, fill=WHITE, border=SLATE_200)
    add_rect(slide, px, card_y, phase_w, Inches(0.05), fill=accent)

    # Title
    add_textbox(slide, px + Inches(0.1), card_y + Inches(0.2), phase_w - Inches(0.2), Inches(0.6),
                title, font_size=13, bold=True, color=SLATE_900, alignment=PP_ALIGN.CENTER)
    # Description
    add_textbox(slide, px + Inches(0.1), card_y + Inches(0.85), phase_w - Inches(0.2), Inches(0.6),
                desc, font_size=10, color=SLATE_500, alignment=PP_ALIGN.CENTER)
    # Date
    add_textbox(slide, px + Inches(0.1), card_y + Inches(1.55), phase_w - Inches(0.2), Inches(0.3),
                date, font_size=10, bold=True, color=accent, alignment=PP_ALIGN.CENTER)

# ── Closing statement ──
add_rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.04), fill=BLUE)
add_textbox(slide, Inches(1), Inches(5.9), Inches(11.3), Inches(1.0),
            "\"From fragmented customer interactions to intelligent,\nprivacy-preserving, AI-driven personalization.\"",
            font_size=26, bold=True, color=SLATE_900, alignment=PP_ALIGN.CENTER)

# ── Speaker notes ──
notes_slide = slide.notes_slide
notes_slide.notes_text_frame.text = (
    "Looking ahead, we have mapped out eight phases of enhancement. "
    "Phase one adds authentication with OAuth 2.0 and SSO. "
    "Phase two brings real-time streaming analytics using Apache Flink for advanced aggregations. "
    "Phase three implements GDPR compliance APIs for data deletion, export, and consent management. "
    "Phase four evolves the platform into a multi-tenant SaaS with organization isolation and usage-based billing. "
    "Phase five enables cloud deployment on AWS, GCP, or Azure using Kubernetes and Terraform. "
    "Phase six introduces predictive customer lifetime value modeling. "
    "Phase seven adds marketing automation with journey builders and campaign triggers. "
    "And phase eight brings an A/B testing engine for multi-variant experimentation. "
    "We believe this platform can transform how enterprises understand and engage their customers — "
    "moving from fragmented interactions to intelligent, privacy-preserving, AI-driven personalization."
)

# ════════════════════════════════════════════════════════════════════════
# SAVE
# ════════════════════════════════════════════════════════════════════════
output_path = "/Users/aman/Projects/CDP/CDP_Hackathon_Presentation.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
print(f"   Slides: {len(prs.slides)}")
print(f"   Size: {math.ceil(prs.slide_width / 914400 * 25.4)}mm x {math.ceil(prs.slide_height / 914400 * 25.4)}mm (widescreen)")
